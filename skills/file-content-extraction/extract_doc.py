#!/usr/bin/env python3
"""
extract_doc.py — Extract and structure content from files and folders.
Supports: PDF, DOCX, PPTX, images, text/code files, and directories.

Usage:
    python extract_doc.py <input_path> [output_dir]
    python extract_doc.py ./my_project [output_dir]     # folder mode

Output:
    output_dir/
    ├── content.md              # Structured extraction (main deliverable)
    └── figures/                # Extracted embedded images
"""

import os
import sys
import subprocess
import shutil
import tempfile
from pathlib import Path


# ---------------------------------------------------------------------------
# Dependency management
# ---------------------------------------------------------------------------

def ensure_dep(package, import_name=None):
    try:
        __import__(import_name or package)
    except ImportError:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", package, "--break-system-packages", "-q"],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )


def ensure_deps_for(ext):
    base = ["Pillow"]
    ext_deps = {
        ".pdf": ["pdfplumber", "pypdf"],
        ".docx": ["python-docx"],
        ".pptx": ["python-pptx"],
    }
    for pkg in base + ext_deps.get(ext, []):
        import_name = {
            "python-docx": "docx", "python-pptx": "pptx",
            "Pillow": "PIL",
        }.get(pkg, pkg)
        ensure_dep(pkg, import_name)


PDFIMAGES_BIN = None

def check_pdfimages():
    """Check if pdfimages is installed (poppler-utils). Try PATH first, then common locations."""
    global PDFIMAGES_BIN
    for candidate in ["pdfimages", "/usr/local/bin/pdfimages", "/opt/homebrew/bin/pdfimages"]:
        try:
            subprocess.run([candidate, "-v"], capture_output=True, check=True)
            PDFIMAGES_BIN = candidate
            return True
        except (FileNotFoundError, subprocess.CalledProcessError):
            continue
    return False


def find_unpack_script():
    """Locate unpack.py from docx or pptx skill."""
    candidates = [
        os.path.expanduser("~/.cursor/skills/docx/scripts/office/unpack.py"),
        os.path.expanduser("~/.cursor/skills/pptx/scripts/office/unpack.py"),
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".svg", ".emf", ".wmf"}
MIN_IMAGE_DIM = 200
MIN_IMAGE_AREA = 80000


# ---------------------------------------------------------------------------
# PDF figure extraction via pdfimages
# ---------------------------------------------------------------------------

TEMPLATE_PAGE_THRESHOLD = 3

def extract_pdf_figures(pdf_path, figures_dir):
    """Extract embedded images from PDF using pdfimages.

    Filters out:
    - Alpha masks (smask type)
    - Small icons/emojis (below MIN_IMAGE_DIM or MIN_IMAGE_AREA)
    - Template/decorative images reused on 3+ pages (backgrounds, logos)

    Returns list of (fig_name, page, width, height).
    Naming: p{page}_fig_{N}.png for traceability.
    """
    if not check_pdfimages():
        print("pdfimages not found -- skipping image extraction. Install: brew install poppler")
        return []

    prefix = os.path.join(figures_dir, "fig")
    subprocess.run([PDFIMAGES_BIN, "-png", pdf_path, prefix],
                    capture_output=True, check=False)

    metadata = _get_pdfimages_metadata(pdf_path)

    obj_page_counts = {}
    for _, _, _, _, obj_id in metadata:
        if obj_id is not None:
            obj_page_counts.setdefault(obj_id, set())
    for page, _, _, _, obj_id in metadata:
        if obj_id is not None and page is not None:
            obj_page_counts[obj_id].add(page)
    template_objs = {oid for oid, pages in obj_page_counts.items()
                     if len(pages) >= TEMPLATE_PAGE_THRESHOLD}

    extracted = sorted(Path(figures_dir).glob("fig-*.png"))
    results = []
    seen_objects = set()
    page_fig_count = {}

    for i, fpath in enumerate(extracted):
        page, width, height, img_type, obj_id = None, 0, 0, "image", None
        if i < len(metadata):
            page, width, height, img_type, obj_id = metadata[i]

        should_remove = False

        if img_type == "smask":
            should_remove = True
        elif width < MIN_IMAGE_DIM or height < MIN_IMAGE_DIM:
            should_remove = True
        elif width * height < MIN_IMAGE_AREA:
            should_remove = True
        elif obj_id is not None and obj_id in template_objs:
            should_remove = True
        elif obj_id is not None and obj_id in seen_objects:
            should_remove = True

        if should_remove:
            fpath.unlink()
            continue

        if obj_id is not None:
            seen_objects.add(obj_id)

        page_num = page or 0
        page_fig_count[page_num] = page_fig_count.get(page_num, 0) + 1
        fig_idx = page_fig_count[page_num]

        new_name = f"p{page_num}_fig_{fig_idx}.png"
        new_path = os.path.join(figures_dir, new_name)
        fpath.rename(new_path)
        results.append((new_name, page, width, height))

    return results


def _get_pdfimages_metadata(pdf_path):
    """Parse pdfimages -list output for page, dimensions, type, and object ID.
    Columns: page num type width height color comp bpc enc interp object ID ...
    """
    result = subprocess.run([PDFIMAGES_BIN, "-list", pdf_path],
                            capture_output=True, text=True, check=False)
    entries = []
    for line in result.stdout.strip().split("\n")[2:]:
        parts = line.split()
        if len(parts) >= 12:
            try:
                page = int(parts[0])
                img_type = parts[2]
                width = int(parts[3])
                height = int(parts[4])
                obj_id = parts[10]
                entries.append((page, width, height, img_type, obj_id))
            except (ValueError, IndexError):
                entries.append((None, 0, 0, "image", None))
        elif len(parts) >= 5:
            try:
                page = int(parts[0])
                img_type = parts[2] if len(parts) > 2 else "image"
                width = int(parts[3])
                height = int(parts[4])
                entries.append((page, width, height, img_type, None))
            except (ValueError, IndexError):
                entries.append((None, 0, 0, "image", None))
    return entries


# ---------------------------------------------------------------------------
# Office (DOCX/PPTX) figure extraction via unpack
# ---------------------------------------------------------------------------

def extract_office_figures(file_path, figures_dir, media_subdir):
    """Extract embedded images from DOCX/PPTX by unpacking the ZIP.

    Filters out small icons/decorations using MIN_IMAGE_DIM and MIN_IMAGE_AREA.
    media_subdir: 'word/media' for DOCX, 'ppt/media' for PPTX.
    Returns list of (fig_name, width, height) tuples.
    """
    unpack_script = find_unpack_script()

    with tempfile.TemporaryDirectory() as tmpdir:
        if unpack_script:
            script_dir = os.path.dirname(os.path.dirname(unpack_script))
            subprocess.run(
                [sys.executable, unpack_script, file_path, tmpdir],
                capture_output=True, check=False,
                cwd=script_dir,
            )
        else:
            import zipfile
            try:
                with zipfile.ZipFile(file_path, "r") as zf:
                    zf.extractall(tmpdir)
            except Exception:
                return []

        media_dir = os.path.join(tmpdir, media_subdir)
        if not os.path.isdir(media_dir):
            return []

        results = []
        fig_num = 0
        for fname in sorted(os.listdir(media_dir)):
            ext = os.path.splitext(fname)[1].lower()
            if ext not in IMAGE_EXTS:
                continue

            src = os.path.join(media_dir, fname)
            w, h = 0, 0
            try:
                from PIL import Image
                with Image.open(src) as img:
                    w, h = img.size
                if w < MIN_IMAGE_DIM or h < MIN_IMAGE_DIM:
                    continue
                if w * h < MIN_IMAGE_AREA:
                    continue
            except Exception:
                pass

            fig_num += 1
            out_ext = ext if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"} else ".png"
            new_name = f"fig_{fig_num}{out_ext}"
            dst = os.path.join(figures_dir, new_name)

            if ext == out_ext:
                shutil.copy2(src, dst)
            else:
                try:
                    from PIL import Image
                    Image.open(src).save(dst)
                except Exception:
                    shutil.copy2(src, os.path.join(figures_dir, f"fig_{fig_num}{ext}"))
                    new_name = f"fig_{fig_num}{ext}"

            results.append((new_name, w, h))

    return results


def _is_decorative_table(table, page_text):
    """Return True if a table is likely a slide template/background element.

    Heuristics:
    - >70% of cells are empty or whitespace
    - All non-empty content is already present in the page text (redundant)
    """
    all_cells = []
    for row in table:
        for cell in row:
            all_cells.append(str(cell or "").strip())
    if not all_cells:
        return True
    non_empty = [c for c in all_cells if c]
    if len(non_empty) / len(all_cells) < 0.3:
        return True
    if page_text and all(c in page_text for c in non_empty):
        return True
    return False


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_pdf(pdf_path, output_dir):
    import pdfplumber
    from pypdf import PdfReader

    reader = PdfReader(pdf_path)
    total_pages = len(reader.pages)
    md = []

    meta = reader.metadata
    if meta:
        parts = []
        for key in ["title", "author", "subject"]:
            val = getattr(meta, key, None)
            if val and val.strip():
                parts.append(f"**{key.capitalize()}:** {val.strip()}")
        if parts:
            md.append(" | ".join(parts))

    md.append(f"**Pages:** {total_pages}\n")

    figures_dir = os.path.join(output_dir, "figures")
    os.makedirs(figures_dir, exist_ok=True)

    figures = extract_pdf_figures(pdf_path, figures_dir)

    fig_by_page = {}
    for fig_name, page, w, h in figures:
        fig_by_page.setdefault(page, []).append((fig_name, w, h))

    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            pnum = i + 1
            text = page.extract_text() or ""
            tables = page.extract_tables() or []

            md.append(f"---\n## Page {pnum}\n")

            if text.strip():
                md.append(text.strip())
                md.append("")

            real_table_idx = 0
            for table in tables:
                if not table or not table[0]:
                    continue
                if _is_decorative_table(table, text):
                    continue
                real_table_idx += 1
                md.append(f"### Table {real_table_idx}\n")
                headers = [str(c or "").strip() for c in table[0]]
                md.append("| " + " | ".join(headers) + " |")
                md.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table[1:]:
                    cells = [str(c or "").strip().replace("\n", " ") for c in row]
                    md.append("| " + " | ".join(cells) + " |")
                md.append("")

            if pnum in fig_by_page:
                for fig_name, w, h in fig_by_page[pnum]:
                    md.append(f"![Page {pnum} figure — {w}x{h}px](figures/{fig_name})\n")

    if figures:
        md.append("---\n## Figure Index\n")
        md.append("| Figure | Page | Size | File |")
        md.append("|--------|------|------|------|")
        for fig_name, page, w, h in figures:
            md.append(f"| {fig_name} | {page or '?'} | {w}x{h} | `figures/{fig_name}` |")
        md.append("")

    summary = [
        f"- Pages: {total_pages}",
        f"- Figures extracted: {len(figures)} (filtered from decorative icons/masks)",
    ]
    return "\n".join(md), summary


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_docx(docx_path, output_dir):
    ensure_deps_for(".docx")
    from docx import Document
    from docx.text.paragraph import Paragraph
    from docx.table import Table as DocxTable

    doc = Document(docx_path)
    file_size = os.path.getsize(docx_path)
    md = []

    props = doc.core_properties
    parts = []
    if props.title: parts.append(f"**Title:** {props.title}")
    if props.author: parts.append(f"**Author:** {props.author}")
    if parts: md.append(" | ".join(parts))

    figures_dir = os.path.join(output_dir, "figures")
    os.makedirs(figures_dir, exist_ok=True)
    figures = extract_office_figures(docx_path, figures_dir, "word/media")

    table_idx = 0
    extracted_text_len = 0

    for elem in doc.element.body:
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

        if tag == "p":
            para = Paragraph(elem, doc)
            text = para.text.strip()
            if not text:
                continue
            extracted_text_len += len(text)
            style = para.style.name if para.style else ""
            if "Heading 1" in style or "Title" in style:
                md.append(f"## {text}\n")
            elif "Heading 2" in style:
                md.append(f"### {text}\n")
            elif "Heading 3" in style:
                md.append(f"#### {text}\n")
            elif "List" in style:
                md.append(f"- {text}")
            else:
                md.append(text)
                md.append("")

        elif tag == "tbl":
            table_idx += 1
            table = DocxTable(elem, doc)
            rows_data = []
            for row in table.rows:
                rows_data.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
            if rows_data:
                md.append(f"### Table {table_idx}\n")
                md.append("| " + " | ".join(rows_data[0]) + " |")
                md.append("| " + " | ".join(["---"] * len(rows_data[0])) + " |")
                for row in rows_data[1:]:
                    md.append("| " + " | ".join(row) + " |")
                md.append("")

    if file_size > 50000 and extracted_text_len < 200:
        md.append("\n⚠️ **Warning:** Very little text extracted relative to file size.")
        md.append("This DOCX may contain text boxes, SmartArt, or equations that python-docx cannot extract.")
        md.append("Consider converting to PDF first for better extraction.\n")

    if figures:
        md.append("\n---\n## Figures\n")
        for fig_name, w, h in figures:
            md.append(f"![Document figure — {w}x{h}px](figures/{fig_name})\n")

    summary = [
        f"- Paragraphs: {len(doc.paragraphs)}",
        f"- Tables: {table_idx}",
        f"- Figures extracted: {len(figures)} (filtered from decorative icons)",
    ]
    return "\n".join(md), summary


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx(pptx_path, output_dir):
    ensure_deps_for(".pptx")
    from pptx import Presentation

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    md = []
    md.append(f"**Slides:** {total_slides}\n")

    figures_dir = os.path.join(output_dir, "figures")
    os.makedirs(figures_dir, exist_ok=True)
    figures = extract_office_figures(pptx_path, figures_dir, "ppt/media")

    for i, slide in enumerate(prs.slides):
        snum = i + 1
        md.append(f"---\n## Slide {snum}\n")
        if slide.slide_layout and slide.slide_layout.name:
            md.append(f"*Layout: {slide.slide_layout.name}*\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        is_title = any(
                            r.font.size and r.font.size.pt and r.font.size.pt >= 20
                            for r in para.runs if r.font.size
                        )
                        md.append(f"### {text}\n" if is_title else text)
                md.append("")
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([c.text.strip().replace("\n", " ") for c in row.cells])
                if rows:
                    md.append("| " + " | ".join(rows[0]) + " |")
                    md.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                    for row in rows[1:]:
                        md.append("| " + " | ".join(row) + " |")
                    md.append("")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md.append(f"> **Notes:** {notes}\n")

    if figures:
        md.append("---\n## Figures\n")
        for fig_name, w, h in figures:
            md.append(f"![Slide figure — {w}x{h}px](figures/{fig_name})\n")

    summary = [f"- Slides: {total_slides}", f"- Figures extracted: {len(figures)} (filtered from decorative icons)"]
    return "\n".join(md), summary


# ---------------------------------------------------------------------------
# Image
# ---------------------------------------------------------------------------

def extract_image(img_path, output_dir):
    ensure_dep("Pillow", "PIL")
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size

    figures_dir = os.path.join(output_dir, "figures")
    os.makedirs(figures_dir, exist_ok=True)
    out_name = f"fig_1{os.path.splitext(img_path)[1]}"
    out = os.path.join(figures_dir, out_name)
    img.save(out)

    md = [
        f"**Type:** Image ({img.format}) | **Dimensions:** {w}x{h}\n",
        f"![](figures/{out_name})\n",
    ]
    return "\n".join(md), [f"- Dimensions: {w}x{h}", "- Figures: 1"]


# ---------------------------------------------------------------------------
# Text / Code
# ---------------------------------------------------------------------------

MAX_TEXT_LINES = 500

def extract_text(file_path, output_dir):
    with open(file_path, "r", errors="replace") as f:
        lines = f.readlines()

    total = len(lines)
    truncated = False
    if total > MAX_TEXT_LINES:
        head = lines[:200]
        tail = lines[-100:]
        omitted = total - 300
        content = "".join(head) + f"\n\n... [{omitted} lines omitted] ...\n\n" + "".join(tail)
        truncated = True
    else:
        content = "".join(lines)

    ext = os.path.splitext(file_path)[1].lstrip(".")
    lang = ext if ext in ["py","js","ts","html","css","json","yaml","yml","sql","sh","bash","go","rs","java","cpp","c","rb","swift","kt"] else ""

    md = [
        f"**Lines:** {total}" + (f" (truncated to first 200 + last 100)" if truncated else ""),
        "",
        f"```{lang}",
        content.rstrip(),
        "```",
    ]
    return "\n".join(md), [f"- Lines: {total}" + (" (truncated)" if truncated else "")]


# ---------------------------------------------------------------------------
# Folder
# ---------------------------------------------------------------------------

BINARY_EXTS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
               ".zip", ".tar", ".gz", ".exe", ".dll", ".so", ".o", ".pyc", ".class",
               ".mp3", ".mp4", ".wav", ".mov", ".avi", ".db", ".sqlite", ".ico", ".woff", ".ttf"}

IGNORE_DIRS = {".git", "node_modules", "__pycache__", ".venv", "venv", ".next", "dist",
               "build", ".cache", ".idea", ".vscode", "vendor", "target", ".DS_Store"}

MAX_FOLDER_FILES = 50


def build_tree(root, max_depth=3, prefix=""):
    entries = sorted(os.listdir(root))
    dirs = [e for e in entries if os.path.isdir(os.path.join(root, e)) and e not in IGNORE_DIRS]
    files = [e for e in entries if os.path.isfile(os.path.join(root, e))]
    lines = []
    for f in files:
        size = os.path.getsize(os.path.join(root, f))
        size_str = f"{size}" if size < 1024 else f"{size/1024:.0f}K" if size < 1048576 else f"{size/1048576:.1f}M"
        lines.append(f"{prefix}├── {f} ({size_str})")
    for d in dirs:
        lines.append(f"{prefix}├── {d}/")
        if max_depth > 1:
            lines.extend(build_tree(os.path.join(root, d), max_depth - 1, prefix + "│   "))
    return lines


def collect_files(root, max_depth=3, _depth=0):
    results = []
    if _depth > max_depth:
        return results
    try:
        entries = sorted(os.listdir(root))
    except PermissionError:
        return results
    for e in entries:
        full = os.path.join(root, e)
        if os.path.isdir(full):
            if e not in IGNORE_DIRS:
                results.extend(collect_files(full, max_depth, _depth + 1))
        elif os.path.isfile(full):
            results.append(full)
    return results


def extract_folder(folder_path, output_dir):
    md = []

    tree = build_tree(folder_path)
    md.append("## File Tree\n")
    md.append("```")
    md.append(os.path.basename(folder_path) + "/")
    md.extend(tree)
    md.append("```\n")

    all_files = collect_files(folder_path)
    md.append(f"**Total files found:** {len(all_files)}\n")

    extractable = {".pdf", ".docx", ".pptx"}
    text_exts = {".txt",".md",".csv",".json",".xml",".yaml",".yml",".log",
                 ".py",".js",".ts",".tsx",".jsx",".html",".css",".sql",".sh",
                 ".bash",".go",".rs",".java",".cpp",".c",".h",".rb",".swift",
                 ".kt",".toml",".cfg",".ini",".env",".gitignore",".dockerignore",
                 "Makefile","Dockerfile"}
    image_exts = {".png",".jpg",".jpeg",".gif",".webp",".bmp"}

    docs = [f for f in all_files if os.path.splitext(f)[1].lower() in extractable]
    texts = [f for f in all_files if os.path.splitext(f)[1].lower() in text_exts or os.path.basename(f) in text_exts]
    images = [f for f in all_files if os.path.splitext(f)[1].lower() in image_exts]
    others = [f for f in all_files if f not in docs + texts + images]

    to_extract = docs + texts[:MAX_FOLDER_FILES - len(docs)]
    skipped = len(texts) + len(docs) - len(to_extract)

    if skipped > 0:
        md.append(f"⚠️ *Extracting {len(to_extract)} files (skipped {skipped} text files to stay within limits).*\n")

    md.append(f"- Documents (PDF/DOCX/PPTX): {len(docs)}")
    md.append(f"- Text/code files: {len(texts)}")
    md.append(f"- Images: {len(images)}")
    md.append(f"- Other (binary/unknown): {len(others)}")
    md.append("")

    summary_parts = [
        f"- Total files: {len(all_files)}",
        f"- Extracted: {len(to_extract)}",
        f"- Images: {len(images)} (listed, not extracted — view individually if needed)",
    ]

    for fpath in to_extract:
        rel = os.path.relpath(fpath, folder_path)
        ext = os.path.splitext(fpath)[1].lower()
        md.append(f"---\n## 📄 {rel}\n")

        try:
            if ext in extractable:
                sub_out = os.path.join(output_dir, os.path.splitext(rel)[0].replace("/", "_"))
                os.makedirs(sub_out, exist_ok=True)
                os.makedirs(os.path.join(sub_out, "figures"), exist_ok=True)
                ensure_deps_for(ext)
                extractors = {".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx}
                content, _ = extractors[ext](fpath, sub_out)
                md.append(content)
            else:
                content, _ = extract_text(fpath, output_dir)
                md.append(content)
        except Exception as e:
            md.append(f"⚠️ Extraction failed: {e}\n")

    if images:
        md.append("---\n## 🖼️ Images (not extracted)\n")
        for img in images:
            rel = os.path.relpath(img, folder_path)
            size = os.path.getsize(img)
            md.append(f"- `{rel}` ({size/1024:.0f}K)")
        md.append("\n*Request individual image extraction if needed.*\n")

    return "\n".join(md), summary_parts


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

FILE_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".png": extract_image, ".jpg": extract_image, ".jpeg": extract_image,
    ".gif": extract_image, ".webp": extract_image, ".bmp": extract_image,
}

TEXT_EXTS = {
    ".txt",".md",".csv",".json",".xml",".yaml",".yml",".log",
    ".py",".js",".ts",".tsx",".jsx",".html",".css",".sql",".sh",
    ".bash",".go",".rs",".java",".cpp",".c",".h",".rb",".swift",
    ".kt",".toml",".cfg",".ini",".env",
}


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_doc.py <input_path> [output_dir]")
        print("  input_path: file or folder")
        sys.exit(1)

    input_path = sys.argv[1]
    if not os.path.exists(input_path):
        print(f"Error: Not found: {input_path}")
        sys.exit(1)

    base = os.path.basename(input_path.rstrip("/"))
    output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.join(
        os.path.dirname(input_path) or ".", f"{base}_extracted"
    )
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(os.path.join(output_dir, "figures"), exist_ok=True)

    is_folder = os.path.isdir(input_path)

    if is_folder:
        md_content, summary_parts = extract_folder(input_path, output_dir)
        title = f"# 📁 {base}/\n"
    else:
        ext = os.path.splitext(input_path)[1].lower()
        if ext in FILE_EXTRACTORS:
            ensure_deps_for(ext)
            md_content, summary_parts = FILE_EXTRACTORS[ext](input_path, output_dir)
        elif ext in TEXT_EXTS:
            md_content, summary_parts = extract_text(input_path, output_dir)
        else:
            print(f"Unsupported file type: {ext}")
            sys.exit(1)
        title = f"# {base}\n"

    header = [title, "## Extraction Summary\n"] + summary_parts + [""]
    final_md = "\n".join(header) + "\n" + md_content

    content_path = os.path.join(output_dir, "content.md")
    with open(content_path, "w") as f:
        f.write(final_md)

    # Clean empty dirs
    for sub in ["figures"]:
        d = os.path.join(output_dir, sub)
        if os.path.isdir(d) and not os.listdir(d):
            os.rmdir(d)

    print(f"\n✅ Extraction complete → {output_dir}/")
    print(f"   content.md ({os.path.getsize(content_path)/1024:.1f}K)")
    for f in sorted(os.listdir(output_dir)):
        fp = os.path.join(output_dir, f)
        if os.path.isfile(fp) and f != "content.md":
            print(f"   {f} ({os.path.getsize(fp)/1024:.1f}K)")
        elif os.path.isdir(fp):
            count = sum(1 for _ in Path(fp).rglob("*") if _.is_file())
            if count:
                print(f"   {f}/ ({count} files)")


if __name__ == "__main__":
    main()
